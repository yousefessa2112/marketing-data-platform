from __future__ import annotations

import sqlite3
from pathlib import Path
from typing import Any, TypedDict, cast

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


PROJECT_ROOT = Path(__file__).resolve().parent.parent
DB_PATH = PROJECT_ROOT / "data" / "marketing.db"
OUTPUT_PATH = PROJECT_ROOT / "dashboard" / "Marketing_Data_Pipeline_Presentation.pptx"

SLIDE_W = 13.333
SLIDE_H = 7.5
MIN_X = 1.0
MIN_Y = 1.5
MAX_TEXT_W = 8.0

HEADER = RGBColor(27, 42, 74)  # #1B2A4A
ACCENT = RGBColor(59, 130, 246)  # #3B82F6
LIGHT = RGBColor(243, 244, 246)  # #F3F4F6
WHITE = RGBColor(255, 255, 255)
DARK = RGBColor(31, 41, 55)


class DataSnapshot(TypedDict):
    channels: list[tuple[str, float, float, float, int, int]]
    row_count: int
    date_range: tuple[str, str, int]
    monthly: list[tuple[str, int]]


def assert_textbox_bounds(x: float, y: float, w: float, h: float) -> None:
    if x < MIN_X:
        raise ValueError(f"Text box left bound too small: {x}")
    if y < MIN_Y:
        raise ValueError(f"Text box top bound too small: {y}")
    if w > MAX_TEXT_W:
        raise ValueError(f"Text box width exceeds limit: {w}")
    if x + w > (SLIDE_W - 1.0):
        raise ValueError("Text box exceeds right margin")
    if y + h > (SLIDE_H - 0.6):
        raise ValueError("Text box exceeds lower slide space")


def add_safe_textbox(slide, x: float, y: float, w: float, h: float):
    assert_textbox_bounds(x, y, w, h)
    return slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))


def style_text(run_font, size: int, color: RGBColor, bold: bool = False) -> None:
    run_font.name = "Calibri"
    run_font.size = Pt(size)
    run_font.bold = bold
    run_font.color.rgb = color


def add_slide_base(prs: Any, title: str, subtitle: str | None = None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE

    # Design elements first (behind text)
    top_bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(SLIDE_W), Inches(1.05))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = HEADER
    top_bar.line.fill.background()

    accent_strip = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(1.05), Inches(SLIDE_W), Inches(0.12)
    )
    accent_strip.fill.solid()
    accent_strip.fill.fore_color.rgb = ACCENT
    accent_strip.line.fill.background()

    title_box = add_safe_textbox(slide, 1.0, 1.5, 8.0, 0.7)
    tf = title_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    style_text(p.font, 30, HEADER, bold=True)

    if subtitle:
        subtitle_box = add_safe_textbox(slide, 1.0, 2.15, 8.0, 0.5)
        stf = subtitle_box.text_frame
        stf.clear()
        sp = stf.paragraphs[0]
        sp.text = subtitle
        style_text(sp.font, 15, ACCENT, bold=False)
    return slide


def add_bullets(slide, items: list[str], x: float = 1.0, y: float = 2.75, w: float = 8.0, h: float = 3.9) -> None:
    if len(items) > 5:
        raise ValueError("No more than 5 bullets per slide")
    box = add_safe_textbox(slide, x, y, w, h)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for i, item in enumerate(items):
        if len(item.split()) > 15:
            raise ValueError(f"Bullet too long: {item}")
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        style_text(p.font, 20, DARK)
        p.space_after = Pt(10)


def query_data() -> DataSnapshot:
    with sqlite3.connect(DB_PATH) as conn:
        raw_channels = conn.execute(
            """
            SELECT campaign_id, ROUND(SUM(cost),2), ROUND(AVG(ctr)*100,2), ROUND(AVG(cpa),2),
                   SUM(clicks), SUM(conversions)
            FROM campaign_performance
            GROUP BY campaign_id
            ORDER BY SUM(cost) DESC
            """
        ).fetchall()
        row_count_raw = conn.execute("SELECT COUNT(*) FROM campaign_performance").fetchone()
        date_range_raw = conn.execute(
            "SELECT MIN(date), MAX(date), COUNT(DISTINCT date) FROM campaign_performance"
        ).fetchone()
        raw_monthly = conn.execute(
            """
            SELECT substr(date,1,7), SUM(clicks)
            FROM campaign_performance
            GROUP BY substr(date,1,7)
            ORDER BY substr(date,1,7)
            """
        ).fetchall()
    channels: list[tuple[str, float, float, float, int, int]] = []
    for row in raw_channels:
        channels.append((str(row[0]), float(row[1]), float(row[2]), float(row[3]), int(row[4]), int(row[5])))
    monthly: list[tuple[str, int]] = []
    for row in raw_monthly:
        monthly.append((str(row[0]), int(row[1])))

    date_range = (
        str(date_range_raw[0]),
        str(date_range_raw[1]),
        int(date_range_raw[2]),
    )
    row_count = int(row_count_raw[0])
    return {"channels": channels, "row_count": row_count, "date_range": date_range, "monthly": monthly}


def add_card(slide, x: float, y: float, w: float, h: float, title: str, value: str, subtitle: str) -> None:
    card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    card.fill.solid()
    card.fill.fore_color.rgb = LIGHT
    card.line.color.rgb = ACCENT
    card.line.width = Pt(1.25)

    title_box = add_safe_textbox(slide, x + 0.2, y + 0.15, min(w - 0.4, 8.0), 0.35)
    tp = title_box.text_frame.paragraphs[0]
    tp.text = title
    style_text(tp.font, 15, HEADER, True)

    value_box = add_safe_textbox(slide, x + 0.2, y + 0.55, min(w - 0.4, 8.0), 0.65)
    vp = value_box.text_frame.paragraphs[0]
    vp.text = value
    style_text(vp.font, 40, ACCENT, True)

    sub_box = add_safe_textbox(slide, x + 0.2, y + 1.25, min(w - 0.4, 8.0), 0.35)
    sp = sub_box.text_frame.paragraphs[0]
    sp.text = subtitle
    style_text(sp.font, 14, DARK, False)


def make_presentation() -> None:
    data = query_data()
    channels = cast(list[tuple[str, float, float, float, int, int]], data["channels"])
    row_count = cast(int, data["row_count"])
    monthly = cast(list[tuple[str, int]], data["monthly"])

    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    # 1. Title
    s1 = add_slide_base(
        prs,
        "Marketing Campaign Performance Analysis",
        "Multi-Channel Insights & Optimization Opportunities",
    )
    hero = s1.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(3.0), Inches(8.0), Inches(2.3))
    hero.fill.solid()
    hero.fill.fore_color.rgb = LIGHT
    hero.line.color.rgb = ACCENT
    title_msg = add_safe_textbox(s1, 1.35, 3.65, 7.3, 1.1)
    tp = title_msg.text_frame.paragraphs[0]
    tp.text = "Executive View: Spend, Efficiency, and Growth Opportunities"
    tp.alignment = PP_ALIGN.CENTER
    style_text(tp.font, 22, HEADER, True)

    # 2. Challenge
    s2 = add_slide_base(prs, "The Challenge")
    add_bullets(
        s2,
        [
            "Marketing data lives in separate channel reports.",
            "Leaders cannot compare performance quickly.",
            "Budget decisions become slower and riskier.",
            "A unified view reveals where money works best.",
        ],
    )

    # 3. Data overview
    s3 = add_slide_base(prs, "Data Overview")
    add_card(s3, 1.0, 2.8, 2.5, 1.8, "Channels", "5", "Search, Social, Display, Email")
    add_card(s3, 3.8, 2.8, 2.5, 1.8, "Daily Records", f"{row_count:,}", "Consistent daily performance tracking")
    add_card(s3, 6.5, 2.8, 2.5, 1.8, "Campaign Window", "6 Months", "180 daily snapshots")

    # 4. Metrics explained
    s4 = add_slide_base(prs, "Key Metrics Explained")
    add_card(s4, 1.0, 2.8, 2.5, 2.2, "CTR", "Clicks / Views", "How many people click")
    add_card(s4, 3.8, 2.8, 2.5, 2.2, "CPC", "Cost / Click", "What each click costs")
    add_card(s4, 6.5, 2.8, 2.5, 2.2, "CPA", "Cost / Customer", "What each new customer costs")

    # 5. Top findings - spend
    s5 = add_slide_base(prs, "Top Findings - Spend")
    add_bullets(s5, ["Search_Google is the biggest spender at $164K."], y=2.75, h=0.7)
    max_spend = max(float(c[1]) for c in channels)
    start_y = 3.45
    for i, channel in enumerate(channels):
        name = str(channel[0]).replace("_", " ")
        spend = float(channel[1])
        y = start_y + i * 0.62
        label = add_safe_textbox(s5, 1.0, y, 2.2, 0.35)
        lp = label.text_frame.paragraphs[0]
        lp.text = name
        style_text(lp.font, 18, DARK)

        bar_w = 4.7 * (spend / max_spend)
        bar = s5.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(3.25), Inches(y + 0.05), Inches(bar_w), Inches(0.28))
        bar.fill.solid()
        bar.fill.fore_color.rgb = ACCENT
        bar.line.fill.background()

        val = add_safe_textbox(s5, 8.1, y, 0.9, 0.35)
        vp = val.text_frame.paragraphs[0]
        vp.text = f"${spend/1000:.1f}K"
        style_text(vp.font, 18, HEADER, True)

    # 6. Top findings - efficiency
    s6 = add_slide_base(prs, "Top Findings - Efficiency")
    left_card = s6.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(2.9), Inches(3.8), Inches(2.1))
    left_card.fill.solid()
    left_card.fill.fore_color.rgb = LIGHT
    left_card.line.color.rgb = ACCENT
    right_card = s6.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(5.2), Inches(2.9), Inches(3.8), Inches(2.1))
    right_card.fill.solid()
    right_card.fill.fore_color.rgb = LIGHT
    right_card.line.color.rgb = ACCENT

    t1 = add_safe_textbox(s6, 1.25, 3.15, 3.3, 0.5)
    t1p = t1.text_frame.paragraphs[0]
    t1p.text = "Most Cost-Efficient"
    style_text(t1p.font, 16, HEADER, True)
    v1 = add_safe_textbox(s6, 1.25, 3.62, 3.3, 0.7)
    v1p = v1.text_frame.paragraphs[0]
    v1p.text = "Email: $7.56 CPA"
    style_text(v1p.font, 38, ACCENT, True)

    t2 = add_safe_textbox(s6, 5.45, 3.15, 3.3, 0.5)
    t2p = t2.text_frame.paragraphs[0]
    t2p.text = "Highest Engagement"
    style_text(t2p.font, 16, HEADER, True)
    v2 = add_safe_textbox(s6, 5.45, 3.62, 3.3, 0.7)
    v2p = v2.text_frame.paragraphs[0]
    v2p.text = "Search: 3.96% CTR"
    style_text(v2p.font, 38, ACCENT, True)

    # 7. Channel comparison
    s7 = add_slide_base(prs, "Channel Comparison")
    cols = [("Channel", 2.8), ("Spend", 1.7), ("CTR", 1.6), ("CPA", 1.9)]
    x = 1.0
    for header, width in cols:
        head = s7.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(2.8), Inches(width), Inches(0.55))
        head.fill.solid()
        head.fill.fore_color.rgb = HEADER
        head.line.color.rgb = WHITE
        htxt = add_safe_textbox(s7, x + 0.08, 2.9, min(width - 0.16, 8.0), 0.35)
        hp = htxt.text_frame.paragraphs[0]
        hp.text = header
        style_text(hp.font, 18, WHITE, True)
        x += width

    for i, row in enumerate(channels):
        y = 3.35 + i * 0.58
        vals = [
            str(row[0]).replace("_", " "),
            f"${float(row[1])/1000:.1f}K",
            f"{float(row[2]):.2f}%",
            f"${float(row[3]):.2f}",
        ]
        x = 1.0
        for j, (_, width) in enumerate(cols):
            cell = s7.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(y), Inches(width), Inches(0.55))
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT if i % 2 == 0 else WHITE
            cell.line.color.rgb = RGBColor(220, 223, 228)
            ctxt = add_safe_textbox(s7, x + 0.08, y + 0.09, min(width - 0.16, 8.0), 0.35)
            cp = ctxt.text_frame.paragraphs[0]
            cp.text = vals[j]
            style_text(cp.font, 18, DARK)
            x += width

    # 8. Trends
    s8 = add_slide_base(prs, "Trends: Clicks Over 6 Months")
    add_bullets(
        s8,
        [
            "Clicks rise through winter and peak in March.",
            "April softens, likely seasonal demand shift.",
            "Overall trend stays healthy across channels.",
        ],
        y=2.75,
        h=1.5,
    )
    max_clicks = max(int(m[1]) for m in monthly)
    for i, (month, clicks) in enumerate(monthly[:6]):
        x = 1.0 + i * 1.25
        h = 1.8 * (int(clicks) / max_clicks) + 0.3
        bar = s8.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(5.15 - h), Inches(0.8), Inches(h))
        bar.fill.solid()
        bar.fill.fore_color.rgb = ACCENT if i == 5 else HEADER
        bar.line.fill.background()
        lbl = add_safe_textbox(s8, x, 5.25, 0.8, 0.32)
        lp = lbl.text_frame.paragraphs[0]
        lp.text = month[5:]
        lp.alignment = PP_ALIGN.CENTER
        style_text(lp.font, 14, DARK, True)

    # 9. Recommendations
    s9 = add_slide_base(prs, "Recommendations")
    add_bullets(
        s9,
        [
            "Increase budget in Search for scalable traffic.",
            "Protect Email investment for efficient conversions.",
            "Improve Display creative and audience targeting.",
            "Review channel mix monthly for seasonality shifts.",
            "Set CPA guardrails before expanding spend.",
        ],
    )

    # 10. How it works
    s10 = add_slide_base(prs, "How It Works")
    steps = ["Collect", "Clean", "Store", "Analyze", "Visualize"]
    for i, step in enumerate(steps):
        x = 1.0 + i * 1.55
        box = s10.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(3.3), Inches(1.3), Inches(1.0))
        box.fill.solid()
        box.fill.fore_color.rgb = LIGHT
        box.line.color.rgb = ACCENT
        txt = add_safe_textbox(s10, x + 0.08, 3.62, 1.14, 0.35)
        tp = txt.text_frame.paragraphs[0]
        tp.text = step
        tp.alignment = PP_ALIGN.CENTER
        style_text(tp.font, 18, HEADER, True)
        if i < 4:
            arr = s10.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.CHEVRON, Inches(x + 1.35), Inches(3.65), Inches(0.18), Inches(0.3))
            arr.fill.solid()
            arr.fill.fore_color.rgb = ACCENT
            arr.line.fill.background()
    add_bullets(s10, ["Automated, repeatable, and ready for regular reporting."], y=4.8, h=0.8)

    # 11. Tools and tech
    s11 = add_slide_base(prs, "Tools & Tech")
    tools = [("Python", "PY"), ("SQL", "SQL"), ("PySpark", "SP"), ("Plotly", "PL"), ("SQLite", "DB"), ("Azure-ready", "AZ")]
    for i, (name, initials) in enumerate(tools):
        tool_row = i // 3
        col = i % 3
        x = 1.0 + col * 2.7
        y = 2.9 + tool_row * 1.8
        circle = s11.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x), Inches(y), Inches(1.0), Inches(1.0))
        circle.fill.solid()
        circle.fill.fore_color.rgb = ACCENT if i % 2 == 0 else HEADER
        circle.line.fill.background()
        itxt = add_safe_textbox(s11, x, y + 0.28, 1.0, 0.4)
        ip = itxt.text_frame.paragraphs[0]
        ip.text = initials
        ip.alignment = PP_ALIGN.CENTER
        style_text(ip.font, 18, WHITE, True)
        nbox = add_safe_textbox(s11, x + 1.1, y + 0.3, 1.6, 0.4)
        np = nbox.text_frame.paragraphs[0]
        np.text = name
        style_text(np.font, 20, DARK, True)

    # 12. Thank you
    s12 = add_slide_base(prs, "Thank You", "Questions & Discussion")
    final_card = s12.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(3.0), Inches(8.0), Inches(1.9))
    final_card.fill.solid()
    final_card.fill.fore_color.rgb = LIGHT
    final_card.line.color.rgb = ACCENT
    thanks = add_safe_textbox(s12, 1.0, 3.55, 8.0, 0.7)
    tp = thanks.text_frame.paragraphs[0]
    tp.text = "Thank you for your time."
    tp.alignment = PP_ALIGN.CENTER
    style_text(tp.font, 32, HEADER, True)

    OUTPUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT_PATH))
    print(f"Presentation generated at: {OUTPUT_PATH}")


if __name__ == "__main__":
    make_presentation()